from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to add a title slide
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text

    # Helper to add a bullet point slide
    def add_bullet_slide(title_text, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = title_text
        tf = body.text_frame
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)

    # 1. Title Slide
    add_title_slide("AgriSense", "Smart Agriculture Monitoring & Disease Detection\nTeam Nerula")

    # 2. Problem Statement
    add_bullet_slide("The Problem", [
        "Delayed Disease Detection: Farmers often notice issues too late.",
        "Manual Inspection: Checking large fields is labor-intensive.",
        "Crop Loss: Significant yield reduction due to untreated pests."
    ])

    # 3. Existing Solutions vs Gap
    add_bullet_slide("Existing Solutions & Gaps", [
        "Current Methods: Manual spotting or lab tests (slow/expensive).",
        "Gap: Lack of real-time, on-field instant analysis.",
        "Barrier: Complex apps that local farmers cannot use easily."
    ])

    # 4. Our Solution
    add_bullet_slide("Our Solution: AgriSense", [
        "AI-Powered Dashboard integrating Drone Feeds.",
        "Instant Disease Detection using Gemini Pro AI.",
        "Simple visual indicators (Red/Green) for farmers.",
        "Works offline and supports local languages."
    ])

    # 5. How It Works
    add_bullet_slide("How It Works", [
        "1. Drone/Camera captures crop image.",
        "2. System analyzes image using local algorithms.",
        "3. Gemini AI validates and identifies specific disease.",
        "4. Actionable solution provided to the farmer instantly."
    ])

    # 6. Technology Stack
    add_bullet_slide("Technology Stack", [
        "Frontend: HTML5, CSS3, JavaScript",
        "AI Engine: Google Gemini Pro API",
        "Mapping: Leaflet JS",
        "Hosting: GitHub Pages"
    ])

    # 7. Key Features
    add_bullet_slide("Key Features", [
        "Live Drone Feed Simulation",
        "Voice Assistant Integration",
        "Multi-language Support (Tamil, Hindi, etc.)",
        "Satellite Disease Mapping"
    ])

    # 8. Use Case
    add_bullet_slide("Use Case Scenario", [
        "Farmer gets an alert for Sector A1.",
        "Opens app, sees red warning zone.",
        "Clicks to analyze: 'Leaf Blast' detected.",
        "App suggests immediate fungicide spray.",
        "Crop saved before spreading."
    ])

    # 9. Impact & Future Scope
    add_bullet_slide("Impact & Future", [
        "Impact: 30% reduction in crop loss.",
        "Future: Integration with IoT soil sensors.",
        "Future: Blockchain for organic certification."
    ])

    # 10. Team & Thank You
    add_bullet_slide("Thank You", [
        "Team Nerula",
        "Surendhar - Full Stack Lead",
        "Contact: https://github.com/8925331076surendhar/agrisense"
    ])

    try:
        prs.save('AgriSense_Presentation.pptx')
        print("Presentation saved successfully as 'AgriSense_Presentation.pptx'")
    except PermissionError:
        print("Error: Could not save file. Please close 'AgriSense_Presentation.pptx' if it is open.")

if __name__ == "__main__":
    try:
        import pptx
        create_presentation()
    except ImportError:
        print("Error: 'python-pptx' library is missing.")
        print("Please run: pip install python-pptx")
